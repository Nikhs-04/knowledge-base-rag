"""
Document processing module for handling multiple file formats
"""
import os
from typing import List, Dict
from pathlib import Path
import pypdf
from docx import Document
from pptx import Presentation


class DocumentProcessor:
    """Process various document formats and extract text"""
    
    SUPPORTED_FORMATS = {'.pdf', '.txt', '.docx', '.pptx'}
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process_document(self, file_path: str) -> Dict[str, any]:
        """
        Process a single document and return structured data
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dict with document metadata and chunks
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if path.suffix.lower() not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {path.suffix}")
        
        # Extract text based on file type
        text = self._extract_text(file_path, path.suffix.lower())
        
        # Create chunks
        chunks = self._create_chunks(text)
        
        return {
            'filename': path.name,
            'file_path': str(path),
            'file_type': path.suffix,
            'total_chunks': len(chunks),
            'chunks': chunks
        }
    
    def _extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from different file formats"""
        
        if file_type == '.pdf':
            return self._extract_pdf(file_path)
        elif file_type == '.txt':
            return self._extract_txt(file_path)
        elif file_type == '.docx':
            return self._extract_docx(file_path)
        elif file_type == '.pptx':
            return self._extract_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    def _create_chunks(self, text: str) -> List[Dict[str, any]]:
        """
        Split text into overlapping chunks
        
        Args:
            text: Full document text
            
        Returns:
            List of chunk dictionaries
        """
        # Clean text
        text = text.strip()
        if not text:
            return []
        
        chunks = []
        start = 0
        chunk_id = 0
        
        while start < len(text):
            # Define chunk boundaries
            end = start + self.chunk_size
            
            # Try to break at sentence/paragraph boundaries
            if end < len(text):
                # Look for sentence endings
                for punct in ['. ', '.\n', '! ', '?\n']:
                    last_break = text.rfind(punct, start, end)
                    if last_break > start + self.chunk_size // 2:
                        end = last_break + 1
                        break
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunks.append({
                    'chunk_id': chunk_id,
                    'text': chunk_text,
                    'start_pos': start,
                    'end_pos': end
                })
                chunk_id += 1
            
            # Move start position with overlap
            start = end - self.chunk_overlap
            
            # Avoid infinite loop
            if start >= len(text):
                break
        
        return chunks
    
    def process_directory(self, directory_path: str) -> List[Dict[str, any]]:
        """
        Process all supported documents in a directory
        
        Args:
            directory_path: Path to directory
            
        Returns:
            List of processed documents
        """
        path = Path(directory_path)
        
        if not path.is_dir():
            raise NotADirectoryError(f"Not a directory: {directory_path}")
        
        processed_docs = []
        
        for file_path in path.rglob('*'):
            if file_path.suffix.lower() in self.SUPPORTED_FORMATS:
                try:
                    doc = self.process_document(str(file_path))
                    processed_docs.append(doc)
                    print(f"✓ Processed: {file_path.name}")
                except Exception as e:
                    print(f"✗ Error processing {file_path.name}: {str(e)}")
        
        return processed_docs